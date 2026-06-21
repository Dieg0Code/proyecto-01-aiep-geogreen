#!/usr/bin/env python
"""
Fusiona en un solo .pptx:
  - slides 1-6 de la base de la directora (su diseño institucional, intacto)
  - slides 7-16 generadas con nuestro sistema (build/parte-geogreen.pptx)

Copia cada slide nuestra (spTree + fondo + imágenes con remapeo de rId) sobre el
layout "Blank" de su base. Salida: GeoGreen-socio-comunitario-2026-06-22.pptx

Uso: uv run --with python-pptx python build/fusionar.py
"""
import copy, os
from io import BytesIO
from pptx import Presentation
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
BASE = os.path.join(HERE, "base", "GEOGREEN-base-directora.pptx")
PART = os.path.join(HERE, "build", "parte-geogreen.pptx")
OUT  = os.path.join(HERE, "GeoGreen-socio-comunitario-2026-06-22.pptx")

R_EMBED = qn("r:embed")
A_BLIP  = qn("a:blip")

base = Presentation(BASE)
part = Presentation(PART)

# 1) Recortar la base a las primeras 5 slides (eliminando sus parts del paquete)
sldIdLst = base.slides._sldIdLst
for sldId in list(sldIdLst)[5:]:
    rId = sldId.get(qn("r:id"))
    if rId:
        base.part.drop_rel(rId)
    sldIdLst.remove(sldId)

# layout "Blank" garantizado (el de la slide 1 de su base)
blank_layout = base.slides[0].slide_layout

def copy_slide(src):
    dst = base.slides.add_slide(blank_layout)
    # limpiar placeholders que add_slide hereda del layout
    for sh in list(dst.shapes):
        sh._element.getparent().remove(sh._element)
    src_cSld = src._element.find(qn("p:cSld"))
    dst_cSld = dst._element.find(qn("p:cSld"))
    # fondo (p:bg) -> primer hijo de cSld
    src_bg = src_cSld.find(qn("p:bg"))
    if src_bg is not None:
        dst_cSld.insert(0, copy.deepcopy(src_bg))
    # formas
    for el in list(src.shapes._spTree):
        if el.tag in (qn("p:nvGrpSpPr"), qn("p:grpSpPr")):
            continue
        dst.shapes._spTree.append(copy.deepcopy(el))
    # remapear imágenes (cross-package): blob fuente -> part destino -> nuevo rId
    for blip in dst._element.iter(A_BLIP):
        rId = blip.get(R_EMBED)
        if not rId:
            continue
        img_part = src.part.related_part(rId)
        _, new_rId = dst.part.get_or_add_image_part(BytesIO(img_part.blob))
        blip.set(R_EMBED, new_rId)
    return dst

for s in part.slides:
    copy_slide(s)

base.save(OUT)
print("OK:", OUT, "·", len(base.slides._sldIdLst), "slides")

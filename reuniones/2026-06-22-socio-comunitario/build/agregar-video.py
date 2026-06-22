#!/usr/bin/env python
"""
Incrusta el video reproducible del modelo 3D en el deck final, sobre la imagen
de portada de la slide "El dispositivo en 3D".

fusionar.py solo copia imagenes (no media), por eso el video se agrega aqui,
despues de la fusion, con python-pptx add_movie (maneja la parte de media + el
poster + el shape de video correctamente). El rectangulo coincide con donde
quedo la imagen de portada en esa slide (mismas pulgadas que en armar-geogreen.js).

Uso: uv run --with python-pptx python build/agregar-video.py
"""
import os
from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DECK = os.path.join(HERE, "GeoGreen-socio-comunitario-2026-06-22.pptx")
VIDEO = os.path.join(HERE, "assets", "render-carcasa-crop.mp4")
POSTER = os.path.join(HERE, "assets", "render-poster.jpg")

# Rectangulo de la imagen de portada en la slide (framed -> inner rect).
LEFT, TOP, WIDTH, HEIGHT = Inches(6.88), Inches(2.11), Inches(5.54), Inches(4.33)
MARCA = "El dispositivo en 3D"

prs = Presentation(DECK)

target = None
for slide in prs.slides:
    for sh in slide.shapes:
        if sh.has_text_frame and MARCA in sh.text_frame.text:
            target = slide
            break
    if target is not None:
        break

if target is None:
    raise SystemExit(f"No se encontro la slide con el titulo '{MARCA}'")

target.shapes.add_movie(
    VIDEO, LEFT, TOP, WIDTH, HEIGHT,
    poster_frame_image=POSTER, mime_type="video/mp4",
)
prs.save(DECK)
print("OK: video incrustado en", os.path.basename(DECK))
